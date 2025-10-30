# src/extract_data.py
import os
import re
from pathlib import Path
from typing import Optional

# Text / PDF / DOCX / PPTX
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Images (OCR)
from PIL import Image, ImageFilter
import pytesseract

# Audio/Video
import ffmpeg  # ffmpeg-python
from pydub import AudioSegment
import speech_recognition as sr

# YouTube
from yt_dlp import YoutubeDL


pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


# -------------------------
# Document extractors
# -------------------------
def extract_from_pdf(file_path: str) -> str:
    try:
        text_parts = []
        reader = PdfReader(file_path)
        for page in reader.pages:
            ptext = page.extract_text()
            if ptext:
                text_parts.append(ptext)
        text = "\n".join(text_parts).strip()
        return text or "[No readable text found in PDF]"
    except Exception as e:
        return f"Error extracting text from PDF: {e}"


def extract_from_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs).strip() or "[No text found in Word file]"
    except Exception as e:
        return f"Error extracting text from Word document: {e}"


def extract_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts).strip() or "[No text found in PPTX]"
    except Exception as e:
        return f"Error extracting text from PPTX: {e}"


def extract_from_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        return f"Error reading text file: {e}"


def extract_from_md(file_path: str) -> str:
    return extract_from_txt(file_path)


# -------------------------
# Image OCR
# -------------------------
def extract_from_image(file_path: str) -> str:
    try:
        img = Image.open(file_path)
        img = img.convert("L").filter(ImageFilter.SHARPEN)
        text = pytesseract.image_to_string(img, config="--psm 6")
        return text.strip() or "[No readable text detected in image]"
    except Exception as e:
        return f"Error extracting text from image: {e}"


# -------------------------
# Audio / Video helpers
# -------------------------
def _convert_to_wav(input_path: str, output_path: Optional[str] = None) -> str:
    """
    Convert input audio/video to WAV (16kHz mono PCM) using ffmpeg-python.
    Returns path to WAV file.
    """
    if output_path is None:
        base = Path(input_path).with_suffix("")
        output_path = str(base) + "_converted.wav"

    try:
        (
            ffmpeg
            .input(input_path)
            .output(output_path, format="wav", acodec="pcm_s16le", ac=1, ar="16000")
            .overwrite_output()
            .run(quiet=True)
        )
        return output_path
    except Exception as e:
        # Try pydub fallback for common audio formats (mp3)
        try:
            audio = AudioSegment.from_file(input_path)
            audio.export(output_path, format="wav")
            return output_path
        except Exception:
            raise RuntimeError(f"FFmpeg/pydub conversion failed: {e}")


def transcribe_wav(wav_path: str) -> str:
    try:
        r = sr.Recognizer()
        with sr.AudioFile(wav_path) as source:
            audio = r.record(source)
        text = r.recognize_google(audio)
        return text
    except Exception as e:
        return f"Error transcribing audio: {e}"


def extract_from_audio(file_path: str) -> str:
    try:
        wav = _convert_to_wav(file_path)
        return transcribe_wav(wav)
    except Exception as e:
        return f"Error extracting text from audio: {e}"


def extract_from_video(file_path: str) -> str:
    try:
        wav = _convert_to_wav(file_path)
        return transcribe_wav(wav)
    except Exception as e:
        return f"Error extracting text from video: {e}"


# -------------------------
# YouTube download + extract
# -------------------------
def download_youtube_audio(url: str, out_template: str = "youtube_audio.%(ext)s") -> str:
    """
    Downloads the best audio and converts to WAV (postprocessor).
    Returns path to WAV file (youtube_audio.wav).
    """
    ydl_opts = {
        'format': 'bestaudio/best',
        'outtmpl': out_template,
        'postprocessors': [{
            'key': 'FFmpegExtractAudio',
            'preferredcodec': 'wav',
            'preferredquality': '192'
        }],
        'quiet': True,
        'no_warnings': True
    }
    try:
        with YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
        # after postprocessor, expected file is youtube_audio.wav
        candidate = "youtube_audio.wav"
        if os.path.exists(candidate):
            return candidate
        # fallback: find file in current directory matching pattern
        base = Path(out_template).stem
        for f in os.listdir("."):
            if re.match(r"youtube_audio.*\.wav", f):
                return f
        # otherwise return expected name
        return candidate
    except Exception as e:
        raise RuntimeError(f"Error downloading YouTube audio: {e}")


def extract_from_youtube(url: str) -> str:
    try:
        wav = download_youtube_audio(url)
        return transcribe_wav(wav)
    except Exception as e:
        return f"Error extracting from YouTube: {e}"


# -------------------------
# Dispatcher
# -------------------------
def extract_text(path_or_url: str) -> str:
    if path_or_url.lower().startswith("http"):
        return extract_from_youtube(path_or_url)

    ext = Path(path_or_url).suffix.lower()
    if ext == ".pdf":
        return extract_from_pdf(path_or_url)
    if ext == ".docx":
        return extract_from_docx(path_or_url)
    if ext == ".pptx":
        return extract_from_pptx(path_or_url)
    if ext in (".txt",):
        return extract_from_txt(path_or_url)
    if ext in (".md",):
        return extract_from_md(path_or_url)
    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
        return extract_from_image(path_or_url)
    if ext in (".mp3", ".wav", ".m4a", ".ogg"):
        return extract_from_audio(path_or_url)
    if ext in (".mp4", ".mov", ".mkv", ".avi"):
        return extract_from_video(path_or_url)

    return f"Unsupported file type: {ext}"
